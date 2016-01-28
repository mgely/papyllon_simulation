import os
import time
import datetime
from utility import byteify,gen_timestamp
from subprocess import Popen,PIPE
from Tkinter import *
import ttk
import json
from qutip import parallel_map
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from shutil import rmtree,copyfile
import signal

class PapyllonSimulation(object):
    """
    - Provides default functions useful in simulations:
        - Spyview compatible files
        - 2D plotting (+background)
        - parralelising functions
    - Provides a way to archive all simulations
        - ask for user to write comments prior and after simulation has run
        - ask user to generate a good png that is representative of the simulation
        - Log everything in a power point       

    """
    def __init__(self,path, settings_list):
        with open(os.path.join(os.path.dirname(__file__),"global.json"),"r") as f:
            self.global_setup =  byteify(json.load(f))

        self.generate_relevant_paths(path)

        self.arg_list = settings_list +\
            ['X_name','X_coord','X_start','X_end','X_points',
            'Y_name','Y_coord','Y_start','Y_end','Y_points',
            'Z_name','Z_coord','Z_start','Z_end','Z_points']


        self.start_commentsgithub = ""
        self.end_comments = ""
        self.load_settings()

    def run(self,parfunc):
        self.ask_for_setup_info()
        self.set_simulation_data_path()
        t0 = time.time()
        self.generate_parameter_space()
        raw_data = parallel_map(parfunc,self.parameter_space, progress_bar = True)
        self.save_dat(raw_data)
        self.generate_spyview_meta()
        p = self.open_spyview()
        t1 = time.time()
        self.generate_timing_info(t0,t1)

        response = self.ask_for_comments()
        if response == -1:
            os.kill(p.pid, signal.SIGINT)
            time.sleep(0.1)
            rmtree(self.simulation_data_path)
        else:
            self.end_comments = response
            self.ask_for_png()
            self.save_meta()
            self.log_on_ppt()
            copyfile(os.path.join(self.manager_path,"settings.json"),
                    os.path.join(self.simulation_data_path,"settings.json"))



    def generate_relevant_paths(self,path):

        # Extract relevant paths
        self.manager_path = os.path.split(__file__)[0]
        self.scripts_path = os.path.split(path)[0]
        self.simulation_type_path = os.path.split(self.scripts_path)[0]
        self.script_name = os.path.split(path)[1].split('.')[0]

        self.settings_file_path = os.path.join(self.manager_path,"settings.json")
        self.spyview_path = self.global_setup["spyview_path"]
        self.log_ppt = os.path.join(self.simulation_type_path,"log.pptx")

    def set_simulation_data_path(self):
        # Generate stamp
        stamp = ""
        stamp += gen_timestamp()
        stamp += "__"

        stamp += self.script_name
        stamp += "__"

        stamp += self.detail

        
        self.simulation_data_path = os.path.join(self.simulation_type_path,"data",stamp)
        os.mkdir(self.simulation_data_path)
        
    def set_parameters(self,p):
        setattr(self, self.X_name, p[0])
        setattr(self, self.Y_name, p[1])
        setattr(self, self.Z_name, p[2])

    def load_settings(self):
        with open(self.settings_file_path,"r") as f:
            settings =  byteify(json.load(f)) # to parse unicodes to strings

        # Apply specified settings
        for arg in self.arg_list:
            try:
                value = settings[arg]
            except:
                raise KeyError('No value for argument '+arg+' found in arg_list')

            setattr(self, arg, value)

    def generate_parameter_space(self):
        self.parameter_space = []
        for X in np.linspace(self.X_start,self.X_end,self.X_points):
            for Y in np.linspace(self.Y_start,self.Y_end,self.Y_points):
                for Z in np.linspace(self.Z_start,self.Z_end,self.Z_points):
                    self.parameter_space.append([X,Y,Z])

    def save_dat(self, raw_data): 
        filename = os.path.join(self.simulation_data_path,
                                "output.dat")
        output = np.c_[self.parameter_space, raw_data]
        print output
        np.savetxt(filename, output)

    def ask_for_comments(self):

        root = Tk()
        root.title("Comments?")
        root.lift()
        mainframe = ttk.Frame(root,padding = (3,3,3,3))
        mainframe.grid(column = 0, row = 0, sticky = (N,W,E,S))

        mainframe.columnconfigure(0,weight = 1)
        mainframe.rowconfigure(0,weight = 1)

        ttk.Label(mainframe,
            text = "Please enter comments for the logging powerpoint..."
            ).grid(column = 0, row = 0, sticky = W)

        comments_entry = Text(mainframe,
            height = 13,
            font = ("Arial", "9"))
        comments_entry.grid(columnspan = 2, row = 1, sticky = (W, E))

        def done():
            self.tmp = str(comments_entry.get("1.0",'end-1c'))
            root.destroy()

        ttk.Button(mainframe,
            text = "Done",
            command = done
            ).grid(column = 0, row = 2,sticky = E)


        def discard():
            self.tmp = -1
            root.destroy()

        ttk.Button(mainframe,
            text = "Discard simulation",
            command = discard
            ).grid(column = 1, row = 2,sticky = W)

        root.mainloop()

        return self.tmp

    def open_spyview(self):
        filename = os.path.join(self.simulation_data_path,"output.dat")
        return Popen([self.spyview_path,filename],stderr = PIPE, stdout=PIPE, shell=False)

    def generate_spyview_meta(self):
        filename = os.path.join(self.simulation_data_path,
                                "output.meta.txt")
        with open(filename,'w') as meta:
            meta.write("#inner loop"+"\n")
            meta.write(str(self.X_points)+"\n")
            meta.write(str(self.X_start)+"\n")
            meta.write(str(self.X_end)+"\n")
            meta.write(self.X_coord+"\n")
            meta.write("#outer loop"+"\n")
            meta.write(str(self.Y_points)+"\n")
            meta.write(str(self.Y_start)+"\n")
            meta.write(str(self.Y_end)+"\n")
            meta.write(self.Y_coord+"\n")
            meta.write("#outer most loop"+"\n")
            meta.write(str(self.Z_points)+"\n")
            meta.write(str(self.Z_start)+"\n")
            meta.write(str(self.Z_end)+"\n")
            meta.write(self.Z_coord+"\n")
            meta.write("#values"+"\n")
            meta.write(str(1)+"\n")
            meta.write("Computed"+"\n")
        pass

    def ask_for_png(self):
        
        root = Tk()
        root.title("Waiting for png...")
        root.lift()
        mainframe = ttk.Frame(root,padding = (3,3,3,3))
        mainframe.grid(column = 0, row = 0, sticky = (N,W,E,S))

        mainframe.columnconfigure(0,weight = 1)
        mainframe.rowconfigure(0,weight = 1)

        ttk.Label(mainframe,
            text = "Please generate a png file for the logging powerpoint..."
            ).grid(column = 0, row = 0, sticky = W)

        ttk.Button(mainframe,
            text = "Done",
            command = root.destroy
            ).grid(column = 0, row = 1)

        def no_image():
            dummy_file = os.path.join(self.manager_path,"dummy.png")
            destination = os.path.join(self.simulation_data_path,"dummy.png")
            copyfile(dummy_file,destination)
            root.destroy()

        ttk.Button(mainframe,
            text = "No image",
            command = no_image
            ).grid(column = 1, row = 1)

        root.mainloop()


    def save_meta(self):
        """
        In a JSON file, store:
            timing information
            settings
            global information from global.JSON
            start and end comments
        """
        filename = os.path.join(self.simulation_data_path,
                                "meta.json")
        with open(filename,"w") as f:
            json.dump(self.__dict__, f, sort_keys=True, indent=4, separators=(',', ': '))

    def log_on_ppt(self):
        # TODO: add all pngs found in file (with spyview api)
        """
        In order of importance, include:
            filename/id
            png (dummy png if no data is found)
            settings
            before-after comment
            timing
        """

        # Check for existance of ppt

        if not os.path.isfile(self.log_ppt):
            template = os.path.join(self.manager_path, "log.pptx")
            destination = os.path.join(self.simulation_type_path, "log.pptx")
            copyfile(template,destination)

        # find png
        for file in os.listdir(self.simulation_data_path):
            if file.endswith(".png") or file.endswith(".PNG"):
                png = os.path.join(self.simulation_data_path, file)
                    
        # load the presentation
        file_already_open = True
        while file_already_open == True:
            try: 
                with open(self.log_ppt,'a') as f:
                    pass
                file_already_open = False
            except IOError, e:
                print str(e)

                root = Tk()
                root.title("Error opening file")
                root.lift()
                mainframe = ttk.Frame(root,padding = (3,3,3,3))
                mainframe.grid(column = 0, row = 0, sticky = (N,W,E,S))

                mainframe.columnconfigure(0,weight = 1)
                mainframe.rowconfigure(0,weight = 1)

                ttk.Label(mainframe,
                    text = "Please close the raw data ppt: "+self.log_ppt
                    ).grid(column = 0, row = 0, sticky = E)

                ttk.Button(mainframe,
                    text = "Done",
                    command = root.destroy
                    ).grid(column = 0, row = 1)
                root.mainloop()

        # add slide
        prs = Presentation(self.log_ppt)
        slide = prs.slides.add_slide(prs.slide_layouts[8])

        title = slide.placeholders[0]
        picture = slide.placeholders[1]
        text = slide.placeholders[2]


        title.text = self.start_comments + '\n' +\
                    self.end_comments


        #render settigs string
        settings = ""
        for arg in self.arg_list:
            settings += arg
            settings += " : "
            settings += str(self.__dict__[arg])
            settings += "\n"

        text.text = self.detail + '\n' +\
                "started: " + self.t_start + "\n" +\
                "ended: " + self.t_stop + "\n" +\
                "lasted: " + self.duration + "\n" +\
                "Settings: \n" +\
                settings


        picture.insert_picture(png)

        # save it
        prs.save(self.log_ppt)

    def ask_for_setup_info(self):
                
        entry_width = 60

        root = Tk()
        root.title("Start simulation")
        root.lift()

        # Put a main frame in the root, which has a padding (NWES)
        mainframe = ttk.Frame(root,padding = (3,3,3,3))

        # Put it in column 0, row 0 of the root and have it stick 
        # to all sides upon resizing
        mainframe.grid(column = 0, row = 0, sticky = (N,W,E,S))

        mainframe.columnconfigure(0,weight = 1)
        mainframe.rowconfigure(0,weight = 1)

        ttk.Label(mainframe,
            text = "Simulation detail: "
            ).grid(column = 0, row = 0, sticky = E)

        detail = StringVar()
        detail_entry = ttk.Entry(mainframe,
            width = entry_width,
            textvariable = detail)
        detail_entry.grid(column = 1, row = 0, sticky = (W, E))



        ttk.Label(mainframe,
            text = "Comments: "
            ).grid(column = 0, row = 1, sticky = E)

        comments_entry = Text(mainframe,
            height = 13,
            font = ("Arial", "9"))
        comments_entry.grid(column = 1, row = 1, sticky = (W, E))
        # To get text use: comments = comments_entry.get("1.0",'end-1c')


        ttk.Label(mainframe,
            text = "Data saved in .../simulation_type/data/(time_stamp)_script_detail"
            ).grid(columnspan = 2, row = 2, sticky = W)

        def start():
            self.start_comments = str(comments_entry.get("1.0",'end-1c'))
            self.detail = str(detail.get())
            root.destroy()

        # Button that calls the function
        ttk.Button(mainframe,
            text = "Start",
            command = start
            ).grid(column = 2, row = 3, sticky = W)

        # Configure all the widgets to have the same padding
        for child in mainframe.winfo_children():
            child.grid_configure(padx = 5, pady = 5)

        # Define where the cursor will be initially
        detail_entry.focus()

        # All of the expansion will happen in the middle when you expand
        # the window
        mainframe.columnconfigure(1,weight = 1)

        # Go
        root.mainloop()

    def generate_timing_info(self,t0,t1):

        self.t_start = str(time.strftime('%Y-%m-%d %H:%M:%S', time.localtime(t0)))
        self.t_stop = str(time.strftime('%Y-%m-%d %H:%M:%S', time.localtime(t1)))
        self.duration = str(datetime.timedelta(seconds=t1 - t0))

        print "Started: "+ self.t_start
        print "Ended: "+ self.t_stop
        print "Measurement time: "+ self.duration

